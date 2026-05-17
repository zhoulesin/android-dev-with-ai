#!/usr/bin/env python3
"""生成 AI 赋能 Android 开发 系列博客的演示 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色方案 ──
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # 深藏青
ACCENT = RGBColor(0x00, 0xD2, 0xFF)          # 亮青
ACCENT2 = RGBColor(0xFF, 0x6B, 0x35)         # 橙
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x44)
GREEN = RGBColor(0x00, 0xE6, 0x76)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
RED = RGBColor(0xFF, 0x45, 0x45)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)         # 卡片背景

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    """给 slide 添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=CARD_BG):
    """添加圆角矩形卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = "PingFang SC"
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, line_spacing=1.6):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "PingFang SC"
        p.space_after = Pt(font_size * (line_spacing - 1))
        p.level = 0

    return tf


def add_top_bar(slide):
    """顶部装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_page_number(slide, num, total=14):
    """页码"""
    add_text(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
             f"{num} / {total}", font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def make_section_slide(title, subtitle, slide_num, total):
    """篇章分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    add_top_bar(slide)

    # 大标题
    add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
             title, font_size=48, color=ACCENT, bold=True)

    # 副标题
    if subtitle:
        add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
                 subtitle, font_size=22, color=LIGHT_GRAY)

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.0), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT2
    line.line.fill.background()

    add_page_number(slide, slide_num, total)
    return slide


def make_content_slide(title, bullet_items, slide_num, total, accent_items=None):
    """内容页：标题 + 要点列表。accent_items 是带颜色标记的要点索引"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide)

    # 标题
    add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.9),
             title, font_size=32, color=WHITE, bold=True)

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.25), Inches(3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT2
    line.line.fill.background()

    # 要点
    if accent_items is None:
        accent_items = set()

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(10.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            label, text = item
            run_label = p.add_run()
            run_label.text = f"▸ {label}："
            run_label.font.size = Pt(18)
            run_label.font.color.rgb = ACCENT if i in accent_items else ACCENT2 if label.startswith("关键") else WHITE
            run_label.font.bold = True
            run_label.font.name = "PingFang SC"

            run_text = p.add_run()
            run_text.text = text
            run_text.font.size = Pt(16)
            run_text.font.color.rgb = LIGHT_GRAY
            run_text.font.name = "PingFang SC"
        else:
            p.text = f"▸ {item}"
            p.font.size = Pt(17)
            p.font.color.rgb = LIGHT_GRAY if i > 0 and accent_items and i not in accent_items else WHITE
            p.font.name = "PingFang SC"

        p.space_after = Pt(12)

    add_page_number(slide, slide_num, total)
    return slide


TOTAL = 14

# ════════════════════════════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)

# 大标题
add_text(slide, Inches(1.5), Inches(1.8), Inches(10.5), Inches(1.5),
         "AI 赋能 Android 开发", font_size=56, color=WHITE, bold=True)

# 副标题
add_text(slide, Inches(1.5), Inches(3.4), Inches(10.5), Inches(0.8),
         "从工具选型到全流程落地的最佳实践", font_size=26, color=ACCENT)

# 描述
add_text(slide, Inches(1.5), Inches(4.4), Inches(10.5), Inches(1.0),
         "GPT + Agent 时代，如何系统化地将 AI 引入 Android 开发全流程——\n模型、工具、Skills、MCP、工作流、避坑、度量，一份完整的实践指南。",
         font_size=18, color=LIGHT_GRAY)

# 装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.6), Inches(3), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT2
line.line.fill.background()

add_text(slide, Inches(1.5), Inches(5.9), Inches(10.5), Inches(0.5),
         "2025  ·  11 篇技术博客  ·  真实项目复盘", font_size=14, color=LIGHT_GRAY)

add_page_number(slide, 1, TOTAL)

# ════════════════════════════════════════════════════════════════
# Slide 2: 目录
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.9),
         "目录", font_size=36, color=WHITE, bold=True)

sections = [
    ("基础篇：认识 AI Coding", [
        "第 1 章 · 背景与趋势 —— 范式跃迁与生态全景",
        "第 2 章 · 大模型选型 —— 7 个模型写 Android 代码实测",
        "第 3 章 · Coding 工具对比 —— Cursor / Copilot / Trae 横向测评",
        "第 4 章 · Skills 与 MCP —— AI 编程的真正威力放大器",
    ]),
    ("实战篇：从技能到全流程", [
        "第 5 章 · 高赞 Skills 对比 —— 20 个热门 Skill 实测评分",
        "第 6 章 · Skills 筛选策略 —— 去重框架与定制指南",
        "第 7 章 · 接入项目全流程 —— PRD → Figma → 编码 → 测试",
        "第 8 章 · 常见陷阱 —— 10 个踩坑实录与解决方案",
        "第 9 章 · 效率保障 —— 采纳率从 40% 提到 85% 的方法",
    ]),
    ("进阶篇：精益与持续", [
        "第 10 章 · Token 经济学 —— Skills 越多越好？实测数据说话",
        "第 11 章 · 持续更新 —— 信息渠道、团队知识库、季度回顾",
    ]),
]

y = 1.5
for title, items in sections:
    add_text(slide, Inches(1.5), Inches(y), Inches(10), Inches(0.5),
             title, font_size=22, color=ACCENT, bold=True)
    y += 0.5
    for item in items:
        add_text(slide, Inches(2.2), Inches(y), Inches(9.5), Inches(0.35),
                 item, font_size=15, color=LIGHT_GRAY)
        y += 0.42
    y += 0.3

add_page_number(slide, 2, TOTAL)

# ════════════════════════════════════════════════════════════════
# Slides 3-14: 正文内容
# ════════════════════════════════════════════════════════════════

# --- 基础篇标题页 ---
make_section_slide("基础篇", "认识 AI Coding —— 建立正确的认知框架", 3, TOTAL)

# Ch1
make_content_slide(
    "第 1 章 · 背景与趋势",
    [
        ("范式跃迁", "Copilot 补全 → GPT-4 对话 → Agent 自主执行，AI 从\"帮手\"变成\"同事\""),
        ("Android 优势", "Kotlin 表达力强、Compose 声明式、Gradle 结构化，天然适合 AI 生成"),
        ("角色转变", "开发者从\"代码写手\"转型为\"架构师 + 审查者\"，核心能力是判断力"),
        ("焦虑与机遇", "AI 不是替代你，是替代不用 AI 的人。掌握 AI Coding 的开发者产出 10x"),
        ("实测数据", "一个促销活动页：传统 3 天 → AI 辅助 40 分钟"),
    ],
    4, TOTAL, accent_items={0, 2}
)

# Ch2
make_content_slide(
    "第 2 章 · 大模型选型：7 个模型写 Android 代码实测",
    [
        ("五大测试任务", "Compose 聊天列表、Gradle 冲突解决、单元测试生成、Java→Kotlin 转换、ANR 分析"),
        ("Compose 能力", "Claude Sonnet/Opus 理解 recomposition 和 state hoisting 最精准；GPT-4o 可用但需提示"),
        ("Gradle 短板", "几乎所有模型在多模块依赖冲突上表现不稳，需要注入 build.gradle 上下文"),
        ("多模型组合", "推荐策略：便宜模型(DeeSeek/Qwen)日常编码 + 贵模型(Claude Opus)架构设计"),
        ("关键发现", "不要迷信排行榜。最能理解你项目架构的模型，比最\"聪明\"的模型更重要"),
    ],
    5, TOTAL, accent_items={4}
)

# Ch3
make_content_slide(
    "第 3 章 · Coding 工具对比：Cursor vs Copilot vs Trae",
    [
        ("7 款工具矩阵", "Cursor、Copilot、Windsurf、Trae、Claude Code CLI、Codex CLI、Augment Code"),
        ("IDE vs CLI", "IDE 工具适合日常编码和 UI 工作；CLI 工具(Claude Code)更适合复杂多步任务"),
        ("核心差异", "Agent 能力 → MCP/Skills 支持 → 终端集成 → Android/Gradle 理解深度"),
        ("Android 实测", "多模块项目中 Claude Code 表现最稳（读 Gradle 依赖图更准），Cursor 在 Compose UI 场景更佳"),
        ("推荐组合", "Cursor 日常编码 + Claude Code 处理复杂任务和架构级变更"),
    ],
    6, TOTAL, accent_items={4}
)

# Ch4
make_content_slide(
    "第 4 章 · Skills 与 MCP：AI 编程的真正威力放大器",
    [
        ("Skills 本质", "不是配置文件，是领域知识注入——让 AI 理解你的项目规范、架构约束和编码习惯"),
        ("SKILL.md 设计", "触发词 + 专业规则 + 示例。一个好的 Compose Skill 能让 AI 从不写 LiveData"),
        ("MCP 协议", "连接外部世界的桥梁——Figma MCP（设计→代码）、GitHub MCP（Issue/PR 操作）、数据库 MCP"),
        ("CLAUDE.md", "最重要的单文件。100 行的项目上下文文件，是后续所有对话质量的基石"),
        ("经验之谈", "Skills 和马步一样——写了 CLAUDE.md 和 Compose Rules，代码质量直接翻倍"),
    ],
    7, TOTAL, accent_items={0, 3}
)

# --- 实战篇标题页 ---
make_section_slide("实战篇", "从 Skills 选型到全流程落地 —— 真实项目中的数据与经验", 8, TOTAL)

# Ch5
make_content_slide(
    "第 5 章 · 高赞 Skills 对比：20 个热门 Skill 实测评分",
    [
        ("评测维度", "有用性、Token 成本、维护质量、Android 适配度 —— 四维度 × 20+ Skills"),
        ("Memory 类", "CLAUDE.md 性价比最高；memory-bank 功能强但 Token 贵；spec-driven 适合大项目"),
        ("流程管理类", "TDD Skill 真正提升代码质量；brainstorming + writing-plans 组合对复杂需求价值极高"),
        ("Android 专项", "社区 Android Skills 较少，建议自建：Compose 组件生成器、Hilt 依赖检查、导航图验证"),
        ("核心感悟", "一个写好 CLD 的项目 + 3 个精准 Skill > 装 20 个通用 Skill"),
    ],
    9, TOTAL, accent_items={4}
)

# Ch6
make_content_slide(
    "第 6 章 · Skills 筛选策略：如何做减法",
    [
        ("五步筛选", "必要性测试 → 频率分析 → A/B 对比 → Token 核算 → 维护度检查"),
        ("功能重叠", "3 个 Code Review Skill 怎么选？原则：选最聚焦的，不选最全面的。覆盖度≠有用度"),
        ("自写 vs 社区", "通用流程用社区 Skill（TDD、Debugging），Android 专项必须自建"),
        ("场景分组", "不要全量加载。开发 5 个 Skill、审查 3 个、调试 3 个，按任务切换"),
        ("经验之谈", "30 个 Skill → Token 翻 3 倍，代码质量反而下降。8-15 个是甜点区间"),
    ],
    10, TOTAL, accent_items={2, 4}
)

# Ch7
make_content_slide(
    "第 7 章 · 接入项目全流程：一个功能的真实复盘",
    [
        ("前置配置", "花 20 分钟写 CLAUDE.md。这是整个流程中 ROI 最高的 20 分钟"),
        ("产品阶段", "PRD → AI 提取技术要点 → 架构方案 → 工时预估。3 分钟出骨架"),
        ("UI 阶段", "Figma → 3 轮迭代法：结构→样式→细节。Design Token 注入是关键"),
        ("开发阶段", "TDD 模式：AI 写测试 → AI 写实现 → 验证。ViewModel 一次过，Repository 需两轮"),
        ("测试阶段", "AI 生成 90% 单测 + 80% UI 测试。覆盖率 45% → 82%"),
        ("最终数据", "总耗时 8.5h（传统 ~20h），AI 写了 70% 代码，Bug 率未升高"),
    ],
    11, TOTAL, accent_items={0, 5}
)

# Ch8
make_content_slide(
    "第 8 章 · 常见陷阱：10 个踩坑实录",
    [
        ("迭代破坏", "AI 重构时误删已有功能 → Git diff 注入 + 模块锁定法"),
        ("Figma 还原", "间距永远差 4dp → Design Token 注入 + 分层生成"),
        ("上下文溢出", "50 模块项目 AI '说胡话' → 模块化 Rules + 按需加载"),
        ("Gradle 幻觉", "AI 建议不存在的依赖版本 → version catalog 上下文注入"),
        ("风格混乱", "3 个不同人写的风格 → ktlint + detekt + 详细 Style Rules"),
        ("安全红线", "AI 生成硬编码密钥 → Security Rules + pre-commit secret 扫描"),
    ],
    12, TOTAL, accent_items={0, 1, 5}
)

# Ch9
make_content_slide(
    "第 9 章 · 效率保障：采纳率从 40% 提到 85%",
    [
        ("Prompt 公式", "角色 + 上下文 + 约束 + 示例 + 输出格式。好 Prompt 让输出质量翻倍"),
        ("Rules 三层设计", "项目级(技术栈+架构) → 模块级(依赖+约束) → 功能级(具体上下文+禁止项)"),
        ("三层防线", "AI 自检(拦截 60%) → ktlint+detekt+Lint(拦截 30%) → 人工审查(拦截 10%)"),
        ("度量体系", "代码采纳率、Bug 引入率、吞吐量的 Git 统计方法 + 季度评估模板"),
        ("70/30 法则", "AI 写 70%(样板、测试、UI)，人工写 30%(架构决策、业务逻辑、安全)"),
    ],
    13, TOTAL, accent_items={2, 3}
)

# --- 进阶篇标题页 ---
make_section_slide("进阶篇", "精益与持续 —— 让 AI Coding 可持续地发挥作用", 14, TOTAL)

# Ch10
make_content_slide(
    "第 10 章 · Token 经济学：Skills 越多越好？",
    [
        ("实测数据", "同任务 5 档配置：0 Skills→$0.18，3个→$0.24，8个→$0.34，25个→$0.59，40个→代码质量反降"),
        ("边际递减", "8-15 个 Skills 是甜点区间。超过 20 个，Token 线性增长但采纳率停滞甚至下降"),
        ("关键发现", "40 个 Skills 的代码质量比 8 个还差——AI 被太多指令搞糊涂了"),
        ("按场景加载", "开发 / 审查 / 调试各一组 Skills，不要全量挂载。用条件触发精准匹配"),
        ("省钱三招", "简单任务用便宜模型、定期裁剪低频 Skills、精简 Rules 文件（200 行以内）"),
    ],
    15, TOTAL, accent_items={0, 2}
)

# Ch11
make_content_slide(
    "第 11 章 · 持续更新：不被 FOMO 拖死",
    [
        ("信息食谱", "2 个 Twitter 号 + 1 份周刊 + GitHub Trending。每月 30 分钟就够了，不要日日刷新"),
        ("团队知识库", "Prompt 模板库 + 踩坑日志 + Skills 配置仓库 + 季度复盘。做到可传承可复用"),
        ("新工具冷静期", "新工具/新模型等 2 周再上。让早期用户踩坑，稳定后再评估"),
        ("季度重置", "每季度 2 小时：审计 Skills → 更新 Rules → 检查指标 → 尝试 1-2 个新工具"),
        ("核心心态", "工具服务于目标。稳定性 > 新奇感。用熟悉的工具写出好代码，比用最新的更值得尊重"),
    ],
    16, TOTAL, accent_items={4}
)

# ════════════════════════════════════════════════════════════════
# 结尾页
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)

add_text(slide, Inches(1.5), Inches(2.5), Inches(10.5), Inches(1.0),
         "Thank You", font_size=48, color=ACCENT, bold=True)

add_text(slide, Inches(1.5), Inches(3.8), Inches(10.5), Inches(0.8),
         "AI 不是替代你，是放大你。\n用得好，你就是 10x 开发者。", font_size=22, color=LIGHT_GRAY)

add_text(slide, Inches(1.5), Inches(5.0), Inches(10.5), Inches(0.8),
         "github.com/cozlya/android-dev-with-ai", font_size=16, color=LIGHT_GRAY)

add_page_number(slide, 17, TOTAL)

# ── 保存 ──
output_path = "/Users/zhouxin/Desktop/cozlya/android-dev-with-ai/AI赋能Android开发-最佳实践.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
